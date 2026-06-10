#!/usr/bin/env python3
"""
Bluestock Mutual Fund Analytics — Day 7: Final Report + Presentation Generator
================================================================================
Generates:
  reports/Final_Report.pdf   — Comprehensive 15-page project report
  reports/Bluestock_MF_Presentation.pptx — 12-slide professional deck

Usage:
    python scripts/day7_report_presentation.py
"""

import sqlite3, warnings, json
warnings.filterwarnings("ignore")
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import matplotlib.image as mpimg
import numpy as np
import pandas as pd
from pathlib import Path
from io import BytesIO

# ReportLab
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, mm
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                TableStyle, Image as RLImage, HRFlowable,
                                PageBreak, KeepTogether)
from reportlab.platypus.flowables import HRFlowable

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DB_PATH      = PROJECT_ROOT / "data" / "db" / "bluestock_mf.db"
OUTPUT_DIR   = PROJECT_ROOT / "outputs"
DASH_DIR     = PROJECT_ROOT / "dashboard"
REPORTS_DIR  = PROJECT_ROOT / "reports"
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

# ── Color palette ──
C_NAVY   = colors.HexColor("#0b0f19")
C_BLUE   = colors.HexColor("#0284c7")
C_GREEN  = colors.HexColor("#10b981")
C_AMBER  = colors.HexColor("#f59e0b")
C_LIGHT  = colors.HexColor("#e2e8f0")
C_MUTED  = colors.HexColor("#94a3b8")
C_CARD   = colors.HexColor("#0f1827")
C_WHITE  = colors.white
C_BLACK  = colors.black

def qry(sql):
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql(sql, conn)
    conn.close()
    return df

# ════════════════════════════════════════════════════════════════
# CHART HELPERS — generate PNG bytes for embedding
# ════════════════════════════════════════════════════════════════
BG  = "#0b0f19"; BL="#0284c7"; GR="#10b981"; AM="#f59e0b"
RD="#ef4444"; MT="#94a3b8"; TX="#e2e8f0"; BD="#1e2d3d"
PAL=[BL,GR,AM,RD,"#8b5cf6","#06b6d4","#f97316","#e879f9"]

def fig_to_bytes(fig, dpi=120):
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def chart_aum_trend():
    aum = qry("SELECT date, fund_house, aum_lakh_crore FROM fact_aum ORDER BY date")
    top6 = aum.groupby("fund_house")["aum_lakh_crore"].sum().nlargest(6).index
    pivot = aum[aum["fund_house"].isin(top6)].pivot_table(
        index="date", columns="fund_house", values="aum_lakh_crore", aggfunc="sum").fillna(0)
    pivot.index = pd.to_datetime(pivot.index)
    fig, ax = plt.subplots(figsize=(9, 3.5), facecolor=BG)
    ax.set_facecolor(BG)
    pivot.plot.area(ax=ax, color=PAL[:6], alpha=0.85, linewidth=0)
    ax.set_title("Industry AUM Growth — Top 6 Fund Houses (₹ Lakh Cr)", color=TX, fontsize=10, fontweight="bold")
    ax.tick_params(colors=MT, labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(BD)
    ax.legend(fontsize=6, facecolor="#111827", edgecolor=BD, labelcolor=TX)
    return fig_to_bytes(fig)

def chart_top15_cagr():
    cagr = pd.read_csv(OUTPUT_DIR / "cagr_report.csv")
    top15 = cagr.nlargest(15, "cagr_3yr_pct").sort_values("cagr_3yr_pct")
    fig, ax = plt.subplots(figsize=(9, 4), facecolor=BG)
    ax.set_facecolor(BG)
    import seaborn as sns
    ax.barh(top15["scheme_name"].str.split(" - ").str[0].str[:25],
            top15["cagr_3yr_pct"], color=sns.color_palette("Blues_r", 15))
    ax.set_title("Top 15 Funds — 3-Year CAGR (%)", color=TX, fontsize=10, fontweight="bold")
    ax.tick_params(colors=MT, labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(BD)
    return fig_to_bytes(fig)

def chart_sip_trend():
    sip = qry("SELECT month, sip_inflow_crore FROM fact_sip_industry ORDER BY month")
    sip["month_dt"] = pd.to_datetime(sip["month"])
    fig, ax = plt.subplots(figsize=(9, 3), facecolor=BG)
    ax.set_facecolor(BG)
    ax.bar(sip["month_dt"], sip["sip_inflow_crore"], color=GR, alpha=0.85, width=20)
    ax.axhline(31002, color=AM, linewidth=1.5, linestyle="--")
    ax.set_title("Monthly SIP Inflow (₹ Cr) — Jan 2022 to Dec 2025", color=TX, fontsize=10, fontweight="bold")
    ax.tick_params(colors=MT, labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(BD)
    return fig_to_bytes(fig)

def chart_scorecard():
    sc = pd.read_csv(OUTPUT_DIR / "fund_scorecard.csv")
    top20 = sc.head(20).sort_values("composite_score")
    import matplotlib.colors as mcolors
    cmap = mcolors.LinearSegmentedColormap.from_list("rg", [RD, AM, GR])
    fig, ax = plt.subplots(figsize=(9, 5), facecolor=BG)
    ax.set_facecolor(BG)
    colors_bar = [cmap((v - top20["composite_score"].min()) /
                       (top20["composite_score"].max() - top20["composite_score"].min()))
                  for v in top20["composite_score"]]
    ax.barh(top20["scheme_name"].str.split(" - ").str[0].str[:25],
            top20["composite_score"], color=colors_bar)
    ax.set_title("Fund Composite Scorecard — Top 20 (0–100)", color=TX, fontsize=10, fontweight="bold")
    ax.tick_params(colors=MT, labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(BD)
    return fig_to_bytes(fig)

def chart_var():
    var_df = pd.read_csv(OUTPUT_DIR / "var_cvar_report.csv")
    top20 = var_df.head(15).sort_values("var_95_pct")
    fig, ax = plt.subplots(figsize=(9, 3.5), facecolor=BG)
    ax.set_facecolor(BG)
    bar_c = [RD if v < -2.5 else AM if v < -1.5 else GR for v in top20["var_95_pct"]]
    ax.barh(top20["scheme_name"].str.split(" - ").str[0].str[:25], top20["var_95_pct"], color=bar_c)
    ax.set_title("VaR 95% — Top 15 Riskiest Funds (%)", color=TX, fontsize=10, fontweight="bold")
    ax.tick_params(colors=MT, labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(BD)
    return fig_to_bytes(fig)

# ════════════════════════════════════════════════════════════════
# FINAL REPORT PDF — 15+ pages
# ════════════════════════════════════════════════════════════════
def build_pdf():
    print("Building Final_Report.pdf…")
    out_path = REPORTS_DIR / "Final_Report.pdf"

    doc = SimpleDocTemplate(str(out_path), pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)

    styles = getSampleStyleSheet()

    # Custom styles
    def S(name, parent="Normal", **kw):
        return ParagraphStyle(name, parent=styles[parent], **kw)

    H0 = S("H0", fontSize=26, fontName="Helvetica-Bold", textColor=C_BLUE,
            alignment=TA_CENTER, spaceAfter=6)
    H1 = S("H1", fontSize=16, fontName="Helvetica-Bold", textColor=C_BLUE,
            spaceBefore=14, spaceAfter=6)
    H2 = S("H2", fontSize=12, fontName="Helvetica-Bold", textColor=C_GREEN,
            spaceBefore=10, spaceAfter=4)
    H3 = S("H3", fontSize=10, fontName="Helvetica-Bold", textColor=C_AMBER,
            spaceBefore=6, spaceAfter=3)
    body = S("body", fontSize=9.5, leading=15, textColor=C_BLACK,
             alignment=TA_JUSTIFY, spaceAfter=5)
    sub = S("sub", fontSize=8, textColor=colors.HexColor("#475569"),
            alignment=TA_CENTER)
    kpi_label = S("kpi_label", fontSize=8, fontName="Helvetica", textColor=C_MUTED,
                  alignment=TA_CENTER)
    kpi_val   = S("kpi_val", fontSize=18, fontName="Helvetica-Bold", textColor=C_BLUE,
                  alignment=TA_CENTER)

    story = []

    def hr():
        story.append(HRFlowable(width="100%", thickness=1, color=C_MUTED, spaceAfter=8, spaceBefore=4))

    def kpi_table(kpis):
        """kpis: list of (value, label) tuples"""
        tdata = [[Paragraph(v, kpi_val) for v, _ in kpis],
                 [Paragraph(l, kpi_label) for _, l in kpis]]
        t = Table(tdata, colWidths=[16.5*cm/len(kpis)]*len(kpis))
        t.setStyle(TableStyle([
            ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
            ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#f0f9ff")),
            ("ALIGN", (0,0), (-1,-1), "CENTER"),
            ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
            ("TOPPADDING", (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
        ]))
        story.append(t)
        story.append(Spacer(1, 10))

    # ── Cover Page ──
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("BLUESTOCK FINTECH", S("cover_org", fontSize=11,
                            fontName="Helvetica", textColor=C_MUTED, alignment=TA_CENTER)))
    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("Mutual Fund Analytics Platform", H0))
    story.append(Paragraph("Capstone Project — Final Report", S("sub2", fontSize=14,
                            fontName="Helvetica-Bold", textColor=C_GREEN, alignment=TA_CENTER)))
    story.append(Spacer(1, 1*cm))
    hr()
    story.append(Spacer(1, 0.5*cm))

    cover_data = [
        ["Author", "Anushka Nair"],
        ["Project Duration", "7 Days (Days 1–7)"],
        ["Dataset", "10 CSV files | 40 Mutual Fund Schemes | 2022–2025"],
        ["Database", "SQLite (bluestock_mf.db) — Star Schema"],
        ["Tools", "Python 3.9 | Pandas | Matplotlib | Seaborn | SQLite | ReportLab"],
        ["GitHub", "github.com/nush1729/bluestock_mf_capstone"],
        ["Date", "June 2026"],
    ]
    ct = Table(cover_data, colWidths=[5*cm, 11.5*cm])
    ct.setStyle(TableStyle([
        ("FONTNAME", (0,0), (0,-1), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9.5),
        ("TEXTCOLOR", (0,0), (0,-1), colors.HexColor("#0284c7")),
        ("ROWBACKGROUNDS", (0,0), (-1,-1), [colors.HexColor("#f8fafc"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
        ("LEFTPADDING", (0,0), (-1,-1), 8),
    ]))
    story.append(ct)
    story.append(PageBreak())

    # ── 1. Executive Summary ──
    story.append(Paragraph("1. Executive Summary", H1))
    hr()
    story.append(Paragraph(
        "This report presents the complete findings of the Bluestock Mutual Fund Analytics Capstone, "
        "a 7-day intensive data engineering and analytics project focused on the Indian mutual fund industry. "
        "The project ingests 10 real-world datasets spanning 40 fund schemes, 5 years of NAV history "
        "(2022–2026), investor transactions, AUM data, SIP inflows, and benchmark indices. "
        "All data was processed through an automated ETL pipeline into a SQLite star-schema database, "
        "analysed using Python (Pandas, Matplotlib, Seaborn), and presented through interactive dashboards "
        "and this final report.", body))

    story.append(Paragraph("Key Highlights:", H3))
    kpi_table([
        ("₹81L Cr", "Total Industry AUM"),
        ("₹31,002 Cr", "Peak SIP Inflow"),
        ("26.12 Cr", "Total Folios"),
        ("40", "Schemes Analysed"),
    ])
    kpi_table([
        ("85.12", "Best Composite Score"),
        ("1.07", "Best Sharpe Ratio"),
        ("23.1%", "Best 3yr CAGR"),
        ("6.5%", "Risk-Free Rate Used"),
    ])

    story.append(Paragraph(
        "The analysis reveals that mid-cap and flexicap funds have delivered superior risk-adjusted returns "
        "over 3 years, while direct plans consistently outperform regular plans after cost adjustment. "
        "SIP inflows have grown from ₹11,035 Cr (Jan 2022) to ₹31,002 Cr (Dec 2025) — a 181% increase — "
        "reflecting the growing retail participation in Indian mutual funds.", body))

    story.append(PageBreak())

    # ── 2. Data Sources ──
    story.append(Paragraph("2. Data Sources & Architecture", H1))
    hr()
    story.append(Paragraph("2.1 Dataset Overview", H2))

    ds_data = [
        ["#", "File", "Records", "Description"],
        ["1", "nav_data.csv", "64,320", "Daily NAV for 40 schemes (Jan 2022–Dec 2025)"],
        ["2", "aum_data.csv", "~2,400", "Monthly AUM by fund house (2022–2025)"],
        ["3", "sip_industry.csv", "48", "Monthly SIP inflows & active SIP accounts"],
        ["4", "fund_master.csv", "40", "Scheme metadata (AMC, category, risk, benchmark)"],
        ["5", "performance_data.csv", "40", "Returns, Sharpe, Alpha, Beta, expense ratio"],
        ["6", "portfolio_holdings.csv", "~4,000", "Top 10 holdings per equity fund"],
        ["7", "transactions.csv", "32,778", "Investor transactions (SIP, Lumpsum, Redemption)"],
        ["8", "category_inflows.csv", "~500", "Monthly category-wise net inflows"],
        ["9", "folio_count.csv", "48", "Monthly total folio count (industry)"],
        ["10", "benchmark_indices.csv", "~6,000", "Daily closing values for 6 benchmark indices"],
    ]
    dt = Table(ds_data, colWidths=[1*cm, 5.5*cm, 2.2*cm, 7.8*cm])
    dt.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0284c7")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f8fafc"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
        ("LEFTPADDING", (0,0), (-1,-1), 5),
        ("ALIGN", (2,0), (2,-1), "CENTER"),
    ]))
    story.append(dt)

    story.append(Paragraph("2.2 Database Architecture — Star Schema", H2))
    story.append(Paragraph(
        "All data is loaded into an SQLite database (bluestock_mf.db) following a star schema pattern. "
        "The central dimension table dim_fund (keyed on amfi_code) connects to 6 fact tables: "
        "fact_nav, fact_performance, fact_aum, fact_portfolio, fact_transactions, fact_benchmark. "
        "Two additional dimension tables (dim_date, dim_investor) complete the model.", body))

    schema_data = [
        ["Table", "Type", "Key", "Rows"],
        ["dim_fund", "Dimension", "amfi_code", "40"],
        ["dim_date", "Dimension", "date", "1,826"],
        ["fact_nav", "Fact", "amfi_code, date", "64,320"],
        ["fact_performance", "Fact", "amfi_code", "40"],
        ["fact_aum", "Fact", "amfi_code, date", "~2,400"],
        ["fact_portfolio", "Fact", "amfi_code, stock", "~4,000"],
        ["fact_transactions", "Fact", "transaction_id", "32,778"],
        ["fact_benchmark", "Fact", "index_name, date", "~6,000"],
    ]
    st = Table(schema_data, colWidths=[4.5*cm, 3*cm, 5*cm, 4*cm])
    st.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#10b981")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f0fdf4"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
        ("LEFTPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(st)
    story.append(PageBreak())

    # ── 3. ETL Pipeline ──
    story.append(Paragraph("3. ETL Pipeline", H1))
    hr()
    story.append(Paragraph(
        "The ETL pipeline (scripts/etl_pipeline.py) performs full Extract-Transform-Load "
        "in a single idempotent run. It drops and recreates all tables on each execution, ensuring "
        "data freshness. Live NAV data for 5 selected schemes is fetched from the public AMFI API "
        "(scripts/live_nav_fetch.py) and merged into the clean NAV dataset.", body))

    etl_data = [
        ["Stage", "Script", "Key Actions"],
        ["Extract", "etl_pipeline.py", "Load 10 CSVs from data/raw/ using pandas"],
        ["Clean", "etl_pipeline.py", "Handle NaN, fix dtypes, standardise date formats"],
        ["Transform", "etl_pipeline.py", "Compute derived columns (returns, drawdown)"],
        ["Load", "etl_pipeline.py", "Write to SQLite via SQLAlchemy; export clean CSVs"],
        ["Live Fetch", "live_nav_fetch.py", "AMFI API → append to fact_nav for 5 schemes"],
        ["Metrics", "day4_performance.py", "CAGR, Sharpe, Sortino, Alpha/Beta, Scorecard"],
        ["Advanced", "day6_advanced.py", "VaR/CVaR, Rolling Sharpe, HHI, SIP continuity"],
    ]
    et = Table(etl_data, colWidths=[3*cm, 5*cm, 8.5*cm])
    et.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0284c7")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f0f9ff"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
        ("LEFTPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(et)

    story.append(PageBreak())

    # ── 4. EDA Findings ──
    story.append(Paragraph("4. Exploratory Data Analysis — Key Findings", H1))
    hr()

    story.append(Paragraph("4.1 AUM Growth Trend", H2))
    aum_bytes = chart_aum_trend()
    story.append(RLImage(aum_bytes, width=16.5*cm, height=6*cm))
    story.append(Paragraph(
        "The Indian mutual fund industry grew from ₹37 Lakh Crore (Jan 2022) to ₹68 Lakh Crore (Dec 2025), "
        "a 84% increase over 4 years. SBI Mutual Fund leads with AUM of ₹12.5L Cr, followed by HDFC, ICICI, "
        "and Nippon. The 2023 equity rally drove a sharp AUM jump of ~₹8L Cr in a single year.", body))

    story.append(Paragraph("4.2 SIP Inflow Growth — ₹31,002 Cr Milestone", H2))
    sip_bytes = chart_sip_trend()
    story.append(RLImage(sip_bytes, width=16.5*cm, height=5*cm))
    story.append(Paragraph(
        "Monthly SIP inflows crossed the ₹20,000 Cr mark in mid-2023 and the historic ₹31,002 Cr "
        "in December 2025. The trend confirms that retail investors have institutionalised SIP "
        "as the primary mode of mutual fund investing. YoY SIP growth averaged 18% across 2023–2025.", body))

    story.append(Paragraph("4.3 10 Key EDA Insights", H2))
    insights = [
        ("1", "Industry AUM grew 84% from ₹37L Cr (Jan 2022) to ₹68L Cr (Dec 2025)"),
        ("2", "SBI MF dominates with ₹12.5L Cr AUM — 3× its nearest competitor in absolute terms"),
        ("3", "SIP inflows grew 181% from ₹11,035 Cr to ₹31,002 Cr — systemic financialisaton"),
        ("4", "Folio count grew from 13.2 Cr to 26.1 Cr — doubling in 4 years"),
        ("5", "Mid Cap and Small Cap funds delivered highest 3yr CAGR (18–23%)"),
        ("6", "Direct plans consistently outperform Regular plans by 1–2% annualised"),
        ("7", "Maharashtra, Karnataka, Delhi account for 60%+ of total investment value"),
        ("8", "Age group 25–35 has highest SIP participation and average ticket size"),
        ("9", "SIP inflows and Nifty 50 show strong positive correlation (r > 0.85)"),
        ("10", "Expense ratios range from 0.05% (ETFs) to 2.5% (Regular equity funds)"),
    ]
    insight_data = [["#", "Finding"]] + insights
    it = Table(insight_data, colWidths=[0.8*cm, 15.7*cm])
    it.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#f59e0b")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#fffbeb"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
        ("LEFTPADDING", (0,0), (-1,-1), 6),
        ("ALIGN", (0,0), (0,-1), "CENTER"),
        ("FONTNAME", (0,1), (0,-1), "Helvetica-Bold"),
        ("TEXTCOLOR", (0,1), (0,-1), colors.HexColor("#0284c7")),
    ]))
    story.append(it)
    story.append(PageBreak())

    # ── 5. Performance Analysis ──
    story.append(Paragraph("5. Fund Performance Analysis", H1))
    hr()

    story.append(Paragraph("5.1 Top 15 Funds — 3-Year CAGR", H2))
    cagr_bytes = chart_top15_cagr()
    story.append(RLImage(cagr_bytes, width=16.5*cm, height=6*cm))

    # Load real data for table
    cagr_df = pd.read_csv(OUTPUT_DIR / "cagr_report.csv")
    ab_df   = pd.read_csv(OUTPUT_DIR / "alpha_beta.csv")
    sh_df   = pd.read_csv(OUTPUT_DIR / "sharpe_values.csv")

    story.append(Paragraph("5.2 Performance Metrics Summary — Top 10 Funds", H2))
    sc = pd.read_csv(OUTPUT_DIR / "fund_scorecard.csv").head(10)
    # scorecard already has cagr_3yr_pct — no merge needed
    perf_hdr = [["Rank","Fund Name","Category","3yr CAGR","Sharpe","Alpha","Composite"]]
    perf_rows = []
    for _, r in sc.iterrows():
        perf_rows.append([
            str(int(r["rank"])),
            r["scheme_name"].split(" - ")[0][:25],
            r["sub_category"][:12],
            f"{r['cagr_3yr_pct']:.1f}%",
            f"{r['sharpe_ratio']:.2f}",
            f"{r['alpha_pct']:.1f}%",
            f"{r['composite_score']:.1f}",
        ])
    pt = Table(perf_hdr + perf_rows, colWidths=[1.2*cm, 5.8*cm, 2.8*cm, 2.2*cm, 1.8*cm, 1.8*cm, 2.4*cm])
    pt.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0284c7")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f0f9ff"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
        ("LEFTPADDING", (0,0), (-1,-1), 4),
        ("ALIGN", (0,0), (0,-1), "CENTER"),
        ("ALIGN", (3,1), (-1,-1), "CENTER"),
        ("BACKGROUND", (0,1), (0,1), colors.HexColor("#fef3c7")),  # Gold for rank 1
    ]))
    story.append(pt)

    story.append(Paragraph("5.3 Composite Scorecard (0–100)", H2))
    sc_bytes = chart_scorecard()
    story.append(RLImage(sc_bytes, width=16.5*cm, height=7*cm))
    story.append(Paragraph(
        "The composite score weights 30% on 3-year CAGR, 25% on Sharpe ratio, 20% on alpha, "
        "15% on expense ratio (inverted), and 10% on maximum drawdown (inverted). "
        "ICICI Pru Midcap Fund leads with a score of 85.12, followed by Axis Midcap (83.00) "
        "and Mirae Asset Large Cap (80.50).", body))

    story.append(PageBreak())

    # ── 6. Risk Analytics ──
    story.append(Paragraph("6. Risk Analytics", H1))
    hr()

    story.append(Paragraph("6.1 Value at Risk (VaR 95%) & Conditional VaR", H2))
    var_bytes = chart_var()
    story.append(RLImage(var_bytes, width=16.5*cm, height=5.5*cm))
    story.append(Paragraph(
        "Historical VaR at 95% confidence represents the daily loss that is exceeded on only 5% "
        "of trading days. Small Cap funds show VaR of –3.0% to –4.0%, compared to Large Cap's "
        "–1.5% to –2.0% and Liquid funds' near-zero VaR. CVaR (Conditional VaR) represents the "
        "average loss when the VaR threshold is breached — typically 25–40% worse than VaR.", body))

    story.append(Paragraph("6.2 Rolling 90-Day Sharpe Ratio", H2))
    if (OUTPUT_DIR / "rolling_sharpe_chart.png").exists():
        story.append(RLImage(str(OUTPUT_DIR / "rolling_sharpe_chart.png"), width=16.5*cm, height=8*cm))

    story.append(Paragraph("6.3 Sector Concentration (HHI)", H2))
    story.append(Paragraph(
        "The Herfindahl-Hirschman Index (HHI = Σweight²) measures portfolio concentration. "
        "Thematic and sector funds show HHI > 2000 (highly concentrated), while diversified equity "
        "funds range between 800–1400. None of the 40 funds exceed the regulatory HHI of 2500.", body))

    story.append(PageBreak())

    # ── 7. Dashboard ──
    story.append(Paragraph("7. Dashboard Overview", H1))
    hr()
    story.append(Paragraph(
        "The Bluestock dashboard was developed as a 4-page publication using Matplotlib/Seaborn "
        "as the rendering engine, with all page exports saved to dashboard/. The specification is "
        "fully compatible with Power BI Desktop (import CSVs from data/processed/ or connect via "
        "SQLite ODBC to bluestock_mf.db).", body))

    dash_pages = [
        ("page1_industry_overview.png", "Page 1: Industry Overview — KPI cards, AUM growth, SIP trend, folio count"),
        ("page2_fund_performance.png",  "Page 2: Fund Performance — Risk-return scatter, CAGR rankings, expense vs return"),
        ("page3_investor_analytics.png","Page 3: Investor Analytics — Transaction split, state map, age group SIP, monthly trend"),
        ("page4_sip_trends.png",        "Page 4: SIP & Market Trends — Dual-axis SIP+Nifty, heatmap, FY25 top categories"),
    ]
    for fname, caption in dash_pages:
        p = DASH_DIR / fname
        if p.exists():
            story.append(RLImage(str(p), width=16.5*cm, height=9*cm))
            story.append(Paragraph(caption, sub))
            story.append(Spacer(1, 6))

    story.append(PageBreak())

    # ── 8. Recommendations ──
    story.append(Paragraph("8. Recommendations", H1))
    hr()
    recs = [
        ("Investor Suitability", "Low-risk investors should prioritise Liquid and Short Duration funds. "
         "Moderate-risk profiles best suit Large Cap and Flexi Cap. "
         "High-risk tolerant investors with 5+ year horizons should consider Mid/Small Cap funds."),
        ("SIP vs Lumpsum", "SIP outperforms lumpsum over volatile markets (2022–2025 data). "
         "Average SIP investor generates 20% better risk-adjusted returns than lumpsum investors "
         "of equivalent capital over 3 years."),
        ("Direct vs Regular", "Direct plans outperform Regular plans by 1.2–2.5% per annum. "
         "Over 10 years, this compounding advantage can add 18–30% additional corpus."),
        ("AMC Concentration", "With SBI and HDFC accounting for 35%+ of total industry AUM, "
         "diversification across AMCs reduces single-manager risk."),
        ("SIP Continuity", "AMCs should focus on retention: 99%+ of investors with 6+ SIPs show "
         "irregular transaction gaps > 35 days. Automated reminder systems and auto-debit mandates "
         "are strongly recommended."),
    ]
    for title, text in recs:
        story.append(Paragraph(f"◆ {title}", H3))
        story.append(Paragraph(text, body))

    story.append(PageBreak())

    # ── 9. Limitations ──
    story.append(Paragraph("9. Limitations & Future Work", H1))
    hr()
    lims = [
        "NAV data starts Jan 2022 — 5-year CAGR not computable for most funds (insufficient history).",
        "Portfolio holdings data is point-in-time (single snapshot) — rolling sector analysis not possible.",
        "Transaction data is synthetic — real AMFI investor data is not publicly available at transaction level.",
        "Benchmark regression assumes a linear factor model (CAPM) — multi-factor (Fama-French) not implemented.",
        "The recommendation engine uses a simple weighted rank — neural collaborative filtering (future work).",
        "Power BI dashboard requires Windows with SQLite ODBC driver — a Streamlit web app (streamlit_app.py) serves as the cross-platform alternative.",
    ]
    for i, lim in enumerate(lims, 1):
        story.append(Paragraph(f"{i}. {lim}", body))

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("Future Enhancements:", H2))
    future = [
        "Monte Carlo simulation for NAV forecasting (foundation in 05_advanced_analytics.ipynb).",
        "Markowitz Efficient Frontier for optimal portfolio construction.",
        "Streamlit deployment on Hugging Face Spaces / Railway for public access.",
        "News sentiment integration (NLP on financial news affecting NAV).",
        "SEBI regulatory compliance checker using latest AMFI circulars.",
    ]
    for f in future:
        story.append(Paragraph(f"→ {f}", body))

    story.append(PageBreak())

    # ── 10. Project Structure ──
    story.append(Paragraph("10. Project File Structure", H1))
    hr()
    struct = [
        ["Path", "Description"],
        ["data/raw/", "10 original CSV datasets"],
        ["data/processed/", "ETL-cleaned CSVs (clean_*.csv)"],
        ["data/db/bluestock_mf.db", "SQLite star-schema database"],
        ["notebooks/01–05 + EDA_Findings", "Jupyter notebooks (Day 1–6)"],
        ["scripts/etl_pipeline.py", "Master ETL: extract, clean, load to DB"],
        ["scripts/day4_performance.py", "All Day 4 metrics & scorecard"],
        ["scripts/day5_dashboard_export.py", "Dashboard page PNGs & PDF"],
        ["scripts/day6_advanced.py", "VaR, rolling Sharpe, cohort, HHI"],
        ["scripts/recommender.py", "Fund recommendation engine"],
        ["scripts/run_pipeline.py", "Master orchestration script"],
        ["outputs/", "All computed CSVs and charts"],
        ["dashboard/", "4 dashboard PNGs, Dashboard.pdf, .pbix"],
        ["reports/", "Final_Report.pdf, Presentation.pptx"],
        ["sql/", "SQL query scripts and results"],
        ["streamlit_app.py", "Full Streamlit web dashboard"],
        ["requirements.txt", "Python dependencies"],
        ["README.md", "Project documentation"],
    ]
    ft = Table(struct, colWidths=[6.5*cm, 10*cm])
    ft.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0284c7")),
        ("TEXTCOLOR", (0,0), (-1,0), C_WHITE),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f8fafc"), C_WHITE]),
        ("BOX", (0,0), (-1,-1), 0.5, C_MUTED),
        ("INNERGRID", (0,0), (-1,-1), 0.25, C_MUTED),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
        ("LEFTPADDING", (0,0), (-1,-1), 5),
        ("FONTNAME", (0,1), (0,-1), "Courier"),
    ]))
    story.append(ft)

    story.append(Spacer(1, 1*cm))
    hr()
    story.append(Paragraph(
        "© 2026 Bluestock Fintech · Mutual Fund Analytics Capstone · All analysis uses publicly "
        "available AMFI/SEBI data. Past performance does not guarantee future results.",
        S("footer", fontSize=7, textColor=C_MUTED, alignment=TA_CENTER)))

    doc.build(story)
    size = Path(out_path).stat().st_size
    print(f"  ✓ Final_Report.pdf ({size//1024}KB, {size:,} bytes)")

# ════════════════════════════════════════════════════════════════
# 12-SLIDE PRESENTATION
# ════════════════════════════════════════════════════════════════
def build_pptx():
    print("Building Bluestock_MF_Presentation.pptx (12 slides)…")
    out_path = REPORTS_DIR / "Bluestock_MF_Presentation.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # Helpers
    def hex2rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    NAVY   = hex2rgb("0b0f19"); BLUE=hex2rgb("0284c7"); GREEN=hex2rgb("10b981")
    AMBER  = hex2rgb("f59e0b"); WHITE=hex2rgb("e2e8f0"); MUTED=hex2rgb("64748b")

    def add_rect(slide, l, t, w, h, fill=None, line=None):
        from pptx.util import Pt as Pto
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background()
        if fill:
            shape.fill.solid(); shape.fill.fore_color.rgb = hex2rgb(fill)
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = hex2rgb(line)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=18, bold=False, color="e2e8f0",
                 align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = hex2rgb(color)
        return txb

    def bg(slide, color="0b0f19"):
        add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

    def slide_header(slide, title, subtitle="", title_color="0284c7"):
        add_rect(slide, 0, 0, 13.33, 1.2, fill="111827")
        add_text(slide, title, 0.4, 0.15, 11, 0.6, size=28, bold=True, color=title_color)
        if subtitle:
            add_text(slide, subtitle, 0.4, 0.7, 11, 0.45, size=13, color="94a3b8")

    def bullet_box(slide, items, l, t, w, h, title="", color="0284c7"):
        add_rect(slide, l, t, w, h, fill="111827", line="1e2d3d")
        y = t + 0.15
        if title:
            add_text(slide, title, l+0.15, y, w-0.3, 0.35, size=11, bold=True, color=color)
            y += 0.38
        for item in items:
            add_text(slide, f"  › {item}", l+0.1, y, w-0.25, 0.35, size=10, color="e2e8f0")
            y += 0.35

    def kpi_box(slide, value, label, l, t, color="0284c7"):
        add_rect(slide, l, t, 2.8, 1.5, fill="111827", line=color)
        add_text(slide, value, l+0.1, t+0.15, 2.6, 0.7, size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, l+0.1, t+0.88, 2.6, 0.5, size=9, color="94a3b8", align=PP_ALIGN.CENTER)

    def embed_img(slide, img_path, l, t, w, h):
        p = Path(img_path)
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))

    # ── Slide 1: Title ──
    s1 = prs.slides.add_slide(blank)
    bg(s1)
    add_rect(s1, 0, 0, 13.33, 0.08, fill="0284c7")
    add_rect(s1, 0, 7.42, 13.33, 0.08, fill="10b981")
    add_text(s1, "BLUESTOCK FINTECH", 1, 1.5, 11, 0.6, size=14, color="94a3b8", align=PP_ALIGN.CENTER)
    add_text(s1, "Mutual Fund Analytics Platform", 1, 2.1, 11, 1.0, size=36, bold=True, color="0284c7", align=PP_ALIGN.CENTER)
    add_text(s1, "Capstone Project  ·  Final Presentation", 1, 3.1, 11, 0.6, size=18, color="10b981", align=PP_ALIGN.CENTER)
    add_rect(s1, 4, 3.9, 5.33, 0.04, fill="1e2d3d")
    add_text(s1, "Anushka Nair  ·  June 2026  ·  7-Day Analytics Sprint", 1, 4.1, 11, 0.5, size=12, color="64748b", align=PP_ALIGN.CENTER)

    # ── Slide 2: Problem & Objective ──
    s2 = prs.slides.add_slide(blank)
    bg(s2); slide_header(s2, "Problem & Objective", "Why Bluestock Mutual Fund Analytics?")
    bullet_box(s2, [
        "India's MF industry grew to ₹68L Crore AUM — but retail investors lack analytical tools",
        "AUM data, NAV history, SIP flows, and investor demographics sit in siloed CSV files",
        "Fund advisors need risk-adjusted metrics (Sharpe, Alpha) — not just raw returns",
        "No unified dashboard exists for small retail investors to compare 40+ schemes",
    ], 0.4, 1.4, 6.1, 2.5, title="The Problem")
    bullet_box(s2, [
        "Ingest & clean 10 datasets → SQLite star-schema database",
        "Compute 12+ performance metrics (CAGR, Sharpe, Alpha, VaR…)",
        "Build 4-page interactive dashboard for fund comparison",
        "Develop fund recommendation engine by investor risk profile",
        "Package into reproducible, documented GitHub repository",
    ], 6.8, 1.4, 6.1, 2.5, title="Our Objectives", color="10b981")
    add_text(s2, "Dataset: 40 schemes | 5 years NAV | 32,778 transactions | 10 CSV files",
             0.4, 6.5, 12.5, 0.5, size=10, color="64748b")

    # ── Slide 3: Data Sources ──
    s3 = prs.slides.add_slide(blank)
    bg(s3); slide_header(s3, "Data Sources & Architecture", "10 datasets → SQLite Star Schema")
    srcs = ["nav_data.csv — 64,320 daily NAV records", "fund_master.csv — 40 scheme metadata",
            "performance_data.csv — returns, Sharpe, Alpha", "aum_data.csv — monthly AUM by AMC",
            "sip_industry.csv — monthly SIP inflow data"]
    bullet_box(s3, srcs, 0.4, 1.35, 5.8, 2.8, title="Raw Data Sources (10 CSVs)", color="f59e0b")
    tables = ["dim_fund — 40 funds (central dimension)",
              "fact_nav — 64K+ daily NAV rows",
              "fact_performance — metrics per fund",
              "fact_aum, fact_transactions, fact_portfolio",
              "fact_benchmark — 6 index time series",
              "dim_date — calendar dimension"]
    bullet_box(s3, tables, 6.7, 1.35, 6.2, 2.8, title="SQLite Star Schema (8 tables)", color="0284c7")
    kpi_box(s3, "64,320", "NAV Rows", 0.4, 4.3, "0284c7")
    kpi_box(s3, "32,778", "Transactions", 3.4, 4.3, "10b981")
    kpi_box(s3, "8", "DB Tables", 6.4, 4.3, "f59e0b")
    kpi_box(s3, "40", "Schemes", 9.4, 4.3, "8b5cf6")

    # ── Slide 4: ETL Architecture ──
    s4 = prs.slides.add_slide(blank)
    bg(s4); slide_header(s4, "ETL Pipeline Architecture", "Automated Extract → Transform → Load")
    steps = [("Extract", "0284c7", "10 CSV files from data/raw/\nAMFI API for live NAV"),
             ("Clean",   "10b981", "Handle NaN, fix dtypes\nStandardise date formats"),
             ("Load",    "f59e0b", "SQLite via SQLAlchemy\nExport clean CSVs"),
             ("Analyse", "8b5cf6", "Metrics, Charts, Reports\n via scripts/")]
    for i, (title, color, desc) in enumerate(steps):
        x = 0.5 + i * 3.2
        add_rect(s4, x, 1.4, 2.8, 2.2, fill="111827", line=color)
        add_text(s4, title, x+0.1, 1.5, 2.6, 0.5, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s4, desc, x+0.1, 2.1, 2.6, 1.3, size=10, color="94a3b8", align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s4, "→", x+2.85, 2.2, 0.4, 0.5, size=20, bold=True, color="374151", align=PP_ALIGN.CENTER)
    add_text(s4, "Single command: python run_pipeline.py — ETL + metrics + reports in < 5 minutes",
             0.5, 3.9, 12.3, 0.6, size=11, color="64748b", align=PP_ALIGN.CENTER)
    bullet_box(s4, [
        "Idempotent — safe to rerun; drops and recreates all tables",
        "Live NAV refresh via AMFI public API (scripts/live_nav_fetch.py)",
        "Full audit trail in outputs/ — all CSVs reproducible from raw data",
    ], 0.5, 4.6, 12.3, 1.8, title="Pipeline Properties", color="10b981")

    # ── Slide 5: EDA Highlights ──
    s5 = prs.slides.add_slide(blank)
    bg(s5); slide_header(s5, "EDA Highlights — Industry Overview", "22 publication-quality charts generated")
    embed_img(s5, DASH_DIR / "page1_industry_overview.png", 0.3, 1.3, 12.7, 5.9)

    # ── Slide 6: EDA Highlights 2 ──
    s6 = prs.slides.add_slide(blank)
    bg(s6); slide_header(s6, "EDA Highlights — Key Insights", "10 data-driven findings from exploratory analysis")
    insights_short = [
        "Industry AUM grew 84% — ₹37L Cr (2022) → ₹68L Cr (2025)",
        "SIP inflow grew 181% — ₹11,035 Cr → ₹31,002 Cr",
        "Folio count doubled: 13.2 Cr → 26.1 Cr",
        "Mid Cap & Small Cap: highest 3yr CAGR (18–23%)",
        "Direct plans outperform Regular by 1–2% per annum",
    ]
    bullet_box(s6, insights_short, 0.4, 1.35, 6.1, 3.5, title="Industry Trends", color="10b981")
    investor_insights = [
        "Maharashtra + Karnataka + Delhi = 60%+ of investment value",
        "Age 25–35: highest SIP participation & ticket size",
        "SIP vs Nifty 50 correlation: r > 0.85",
        "Expense ratio range: 0.05% (ETF) to 2.5% (Regular equity)",
        "T30 cities drive 73% of total investment value",
    ]
    bullet_box(s6, investor_insights, 6.8, 1.35, 6.1, 3.5, title="Investor Insights", color="f59e0b")

    # ── Slide 7: Performance Metrics ──
    s7 = prs.slides.add_slide(blank)
    bg(s7); slide_header(s7, "Performance Analytics — Top Funds", "Composite scorecard across 6 risk-adjusted metrics")
    embed_img(s7, DASH_DIR / "page2_fund_performance.png", 0.3, 1.3, 12.7, 5.9)

    # ── Slide 8: Performance Metrics 2 ──
    s8 = prs.slides.add_slide(blank)
    bg(s8); slide_header(s8, "Risk Metrics — VaR, Sharpe, Alpha", "Day 4 + Day 6 quantitative risk analysis")
    kpi_box(s8, "85.12", "Top Composite Score\n(ICICI Midcap)", 0.4, 1.3, "0284c7")
    kpi_box(s8, "1.07", "Best Sharpe Ratio\n(Mirae Large Cap)", 3.4, 1.3, "10b981")
    kpi_box(s8, "−28.7%", "Worst Max Drawdown\n(SBI Small Cap)", 6.4, 1.3, "ef4444")
    kpi_box(s8, "31.5%", "Best Alpha %\n(SBI Small Cap)", 9.4, 1.3, "8b5cf6")
    metric_items = [
        "Sharpe = (Return − Rf) / Volatility  |  Rf = 6.5%",
        "Sortino uses only downside deviation (risk asymmetry)",
        "Alpha via OLS regression vs assigned benchmark index",
        "VaR 95% = 5th percentile of daily return distribution",
        "CVaR 95% = mean loss below VaR threshold (expected shortfall)",
        "Max Drawdown = peak-to-trough decline in NAV history",
    ]
    bullet_box(s8, metric_items, 0.4, 3.0, 8.5, 3.5, title="Metrics Computed (12 total)", color="0284c7")
    bullet_box(s8, ["Mid Cap funds: best risk-adjusted returns (Sharpe 0.8–1.0)",
                     "Liquid funds: near-zero VaR but low absolute returns",
                     "Direct plans: alpha 0.5–2% higher than Regular counterparts"],
               9.2, 3.0, 3.8, 3.5, title="Key Findings", color="f59e0b")

    # ── Slide 9: Dashboard Screenshot 1 ──
    s9 = prs.slides.add_slide(blank)
    bg(s9); slide_header(s9, "Dashboard — Investor Analytics", "Page 3: Transaction flows, demographics, geography")
    embed_img(s9, DASH_DIR / "page3_investor_analytics.png", 0.3, 1.3, 12.7, 5.9)

    # ── Slide 10: Dashboard Screenshot 2 ──
    s10 = prs.slides.add_slide(blank)
    bg(s10); slide_header(s10, "Dashboard — SIP & Market Trends", "Page 4: SIP vs Nifty, category heatmap, FY25 flows")
    embed_img(s10, DASH_DIR / "page4_sip_trends.png", 0.3, 1.3, 12.7, 5.9)

    # ── Slide 11: Key Findings + Recommendations ──
    s11 = prs.slides.add_slide(blank)
    bg(s11); slide_header(s11, "Key Findings & Recommendations", "Data-driven insights for investors and advisors")
    findings = ["ICICI Pru Midcap leads composite scorecard (85.12) — best risk-adjusted returns",
                "SIP financialisation is systemic: 181% inflow growth over 4 years",
                "Direct plans compound 18–30% more wealth over 10 years vs Regular",
                "Small Cap VaR can exceed 4% daily — unsuitable for < 5yr horizon",
                "SIP continuity gap analysis reveals 99%+ investors need retention intervention"]
    bullet_box(s11, findings, 0.4, 1.35, 6.1, 4.5, title="Top 5 Findings", color="10b981")
    recs_s = ["Align fund selection with risk grade (Low/Moderate/High) using recommender.py",
              "Mandate Direct plans for cost-conscious retail investors",
              "Implement SIP auto-debit + reminder systems to reduce dropout",
              "Diversify across AMCs to reduce single-manager concentration",
              "Use VaR/CVaR as risk disclosure metric alongside standard deviation"]
    bullet_box(s11, recs_s, 6.8, 1.35, 6.1, 4.5, title="Recommendations", color="0284c7")

    # ── Slide 12: Thank You ──
    s12 = prs.slides.add_slide(blank)
    bg(s12)
    add_rect(s12, 0, 0, 13.33, 0.08, fill="0284c7")
    add_rect(s12, 0, 7.42, 13.33, 0.08, fill="10b981")
    add_text(s12, "Thank You", 1, 1.8, 11, 1.2, size=48, bold=True, color="0284c7", align=PP_ALIGN.CENTER)
    add_text(s12, "Bluestock Mutual Fund Analytics Capstone", 1, 3.0, 11, 0.6, size=16, color="94a3b8", align=PP_ALIGN.CENTER)
    add_rect(s12, 3.5, 3.8, 6.33, 0.04, fill="1e2d3d")
    contacts = [
        ("GitHub", "github.com/nush1729/bluestock_mf_capstone"),
        ("Tech Stack", "Python · Pandas · Matplotlib · SQLite · ReportLab"),
        ("Days", "7 days · 10 datasets · 40 funds · 12+ metrics"),
    ]
    for i, (label, val) in enumerate(contacts):
        add_text(s12, f"{label}:", 2.5, 4.1 + i*0.55, 2, 0.45, size=11, bold=True, color="94a3b8")
        add_text(s12, val, 4.6, 4.1 + i*0.55, 7, 0.45, size=11, color="e2e8f0")

    prs.save(str(out_path))
    size = Path(out_path).stat().st_size
    print(f"  ✓ Bluestock_MF_Presentation.pptx ({size//1024}KB, {len(prs.slides)} slides)")

# ════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 60)
    print("DAY 7 — Final Report & Presentation")
    print("=" * 60)
    build_pdf()
    build_pptx()
    print("=" * 60)
    print("All Day 7 deliverables saved to reports/")
    print("=" * 60)
