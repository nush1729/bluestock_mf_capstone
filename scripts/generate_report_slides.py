#!/usr/bin/env python3
"""
Bluestock Mutual Fund Analytics — Final Report & Presentation Generator (D7)
==========================================================================
Generates the Final_Report.pdf and Presentation.pptx files.

Usage:
    python scripts/generate_report_slides.py
"""

import sys
import logging
from pathlib import Path
import sqlite3
import pandas as pd

# ReportLab for PDF
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch

# python-pptx for Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logging.basicConfig(level=logging.INFO, format="%(message)s")
logger = logging.getLogger("ReportGen")

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DB_PATH = PROJECT_ROOT / "data" / "db" / "bluestock_mf.db"
OUTPUT_DIR = PROJECT_ROOT / "reports"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Data Fetching
# ---------------------------------------------------------------------------
def fetch_data():
    if not DB_PATH.exists():
        logger.error(f"Database not found at {DB_PATH}")
        sys.exit(1)
        
    conn = sqlite3.connect(DB_PATH)
    
    # Top 5 Funds
    top_funds = pd.read_sql("""
        SELECT f.scheme_name, f.sub_category, f.fund_house, p.return_3yr_pct, p.sharpe_ratio
        FROM dim_fund f
        JOIN fact_performance p ON f.amfi_code = p.amfi_code
        ORDER BY p.return_3yr_pct DESC LIMIT 5
    """, conn)
    
    # Top AMCs
    top_amcs = pd.read_sql("""
        SELECT fund_house, MAX(aum_crore) as aum
        FROM fact_aum GROUP BY fund_house ORDER BY aum DESC LIMIT 5
    """, conn)
    
    # Metrics
    metrics = {
        "Total Funds": pd.read_sql("SELECT COUNT(*) as c FROM dim_fund", conn).iloc[0,0],
        "Avg 3Y Return": pd.read_sql("SELECT AVG(return_3yr_pct) as c FROM fact_performance", conn).iloc[0,0],
        "Total AUM (Lakh Cr)": pd.read_sql("SELECT SUM(aum_lakh_crore) as c FROM fact_aum WHERE date = (SELECT MAX(date) FROM fact_aum)", conn).iloc[0,0],
    }
    
    conn.close()
    return top_funds, top_amcs, metrics

# ---------------------------------------------------------------------------
# PDF Report Generation (ReportLab)
# ---------------------------------------------------------------------------
def generate_pdf(top_funds, top_amcs, metrics):
    pdf_path = OUTPUT_DIR / "Final_Report.pdf"
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = styles['Title']
    title_style.textColor = colors.HexColor("#0056b3")
    h1_style = styles['Heading1']
    h1_style.textColor = colors.HexColor("#0056b3")
    h2_style = styles['Heading2']
    normal_style = styles['Normal']
    
    story = []
    
    # Title Page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("Bluestock Mutual Fund Analytics", title_style))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("Capstone Project Final Report", h1_style))
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("Prepared by: Data Engineering & Analytics Team", normal_style))
    story.append(PageBreak())
    
    # 1. Executive Summary
    story.append(Paragraph("1. Executive Summary", h1_style))
    story.append(Paragraph(
        "This report outlines the end-to-end implementation of the Bluestock Mutual Fund Analytics Platform. "
        "The project consolidates fragmented mutual fund data into a centralized, production-ready SQLite data warehouse. "
        "It includes robust data engineering, automated ETL pipelines, performance & risk analytics, and an interactive Streamlit dashboard.", 
        normal_style
    ))
    story.append(Spacer(1, 0.2*inch))
    
    story.append(Paragraph(f"• Total Funds Tracked: {metrics['Total Funds']}", normal_style))
    story.append(Paragraph(f"• Industry AUM: ₹{metrics['Total AUM (Lakh Cr)']:.2f} Lakh Crore", normal_style))
    story.append(Paragraph(f"• Average 3Y CAGR: {metrics['Avg 3Y Return']:.2f}%", normal_style))
    story.append(Spacer(1, 0.3*inch))
    
    # 2. Architecture & Data Engineering
    story.append(Paragraph("2. Architecture & Data Engineering", h1_style))
    story.append(Paragraph(
        "The system uses a classic ETL architecture. Python and pandas handle extraction from 10 raw CSV datasets. "
        "Transformations include date parsing, missing value forward-filling (ffill for NAVs on weekends/holidays), "
        "and metric computations. The loaded data is structured in a 10-table Star Schema within SQLite.",
        normal_style
    ))
    story.append(Spacer(1, 0.3*inch))
    
    # 3. Performance Analytics
    story.append(Paragraph("3. Top Performing Funds", h1_style))
    
    # Format table data
    table_data = [["Scheme Name", "Category", "AMC", "3Y Return", "Sharpe Ratio"]]
    for _, row in top_funds.iterrows():
        table_data.append([
            row['scheme_name'][:30] + "...", 
            row['sub_category'], 
            row['fund_house'], 
            f"{row['return_3yr_pct']:.2f}%", 
            f"{row['sharpe_ratio']:.2f}"
        ])
        
    t = Table(table_data, colWidths=[2.5*inch, 1.2*inch, 1.5*inch, 1*inch, 1*inch])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#0056b3")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('BACKGROUND', (0,1), (-1,-1), colors.beige),
        ('GRID', (0,0), (-1,-1), 1, colors.black),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.3*inch))
    
    # 4. Advanced Analytics & Machine Learning
    story.append(Paragraph("4. Advanced Analytics & Machine Learning", h1_style))
    story.append(Paragraph(
        "• <b>K-Means Clustering:</b> Grouped 40 funds into 3 clusters based on risk-return profiles.<br/>"
        "• <b>Monte Carlo Simulation:</b> Projected 5-year NAV growth paths using historical volatility.<br/>"
        "• <b>Recommendation Engine:</b> Developed a content-based filtering system matching investor risk appetite to optimal funds.<br/>"
        "• <b>Markowitz Efficient Frontier:</b> Calculated optimal portfolio weights for 5 selected equity funds.",
        normal_style
    ))
    story.append(Spacer(1, 0.3*inch))
    
    # 5. Limitations & Future Scope
    story.append(Paragraph("5. Limitations & Future Scope", h1_style))
    story.append(Paragraph(
        "• <b>Limitations:</b> Static historical data dependency. AUM data is quarterly, limiting daily correlation analysis.<br/>"
        "• <b>Future Scope:</b> Integration of real-time market news sentiment analysis, deployment of the database to PostgreSQL on AWS/GCP, "
        "and building user-specific portfolio tracking in the Streamlit app.",
        normal_style
    ))

    # Build PDF
    doc.build(story)
    logger.info(f"Generated PDF: {pdf_path}")

# ---------------------------------------------------------------------------
# PPTX Presentation Generation
# ---------------------------------------------------------------------------
def generate_pptx(top_funds, top_amcs, metrics):
    pptx_path = OUTPUT_DIR / "Presentation.pptx"
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Bluestock Mutual Fund Analytics"
    subtitle.text = "Capstone Project Presentation\nData Engineering & Analytics Team"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary & KPIs"
    tf = body_shape.text_frame
    tf.text = "Project Objectives Delivered:"
    
    p = tf.add_paragraph()
    p.text = "Automated ETL Pipeline across 10 datasets"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "SQLite Star Schema Data Warehouse"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Analyzed {metrics['Total Funds']} Funds across {len(top_amcs)} AMCs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Industry Average 3Y Return: {metrics['Avg 3Y Return']:.2f}%"
    p.level = 1
    
    # Slide 3: Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "5-Layer Data Architecture:"
    
    p = tf.add_paragraph()
    p.text = "1. Extract: 10 Raw CSVs + mfapi.in Live NAV"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Transform: Pandas cleaning, ffill, daily returns"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Load: SQLite (dim_fund, fact_nav, fact_performance, etc.)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Analyze: Risk Metrics (Sharpe, Alpha, VaR)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Visualise: Interactive Streamlit Dashboard"
    p.level = 1
    
    # Slide 4: Top Funds Table
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    title = slide.shapes.title
    title.text = "Top Performing Funds (3Y Return)"
    
    rows, cols = len(top_funds) + 1, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)
    
    # Headers
    headers = ["Scheme Name", "Category", "3Y Return", "Sharpe Ratio"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    # Data
    for row_idx, row in top_funds.reset_index().iterrows():
        table.cell(row_idx+1, 0).text = row['scheme_name'][:40]
        table.cell(row_idx+1, 1).text = row['sub_category']
        table.cell(row_idx+1, 2).text = f"{row['return_3yr_pct']:.2f}%"
        table.cell(row_idx+1, 3).text = f"{row['sharpe_ratio']:.2f}"
    
    # Slide 5: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion & Recommendations"
    tf = body_shape.text_frame
    tf.text = "Key Takeaways:"
    
    p = tf.add_paragraph()
    p.text = "Significant AUM concentration in top 3 AMCs (SBI, ICICI, HDFC)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SIP Inflows show strong resilience against market volatility"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Small Cap funds demonstrated highest alpha but with elevated Max Drawdown"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Platform is ready for production deployment"
    p.level = 1
    
    prs.save(str(pptx_path))
    logger.info(f"Generated Presentation: {pptx_path}")

def main():
    logger.info("Starting Report and Presentation generation...")
    top_funds, top_amcs, metrics = fetch_data()
    generate_pdf(top_funds, top_amcs, metrics)
    generate_pptx(top_funds, top_amcs, metrics)
    logger.info("Report generation complete.")

if __name__ == "__main__":
    main()
